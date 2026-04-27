"""Helpers para construir decks Educa Edtech a partir de la plantilla.

Uso típico:

    from pathlib import Path
    import sys
    SKILL = Path(__file__).parent.parent   # carpeta del skill
    sys.path.insert(0, str(SKILL / "scripts"))
    from deck_builder import open_template, duplicate_slide, remove_slide, \
                             replace_text_in_slide, insert_png_into_placeholder
    from palette import CARBON, AZUL_EDUCA, FONT_TITLE, FONT_BODY

    prs = open_template()
    # duplicar la portada (slide 2, índice 1) como punto de partida
    portada = duplicate_slide(prs, 1)
    replace_text_in_slide(portada, {"EDUCA EDTECH": "EDUCA EDTECH", "Group": "Group"})
    prs.save("salida.pptx")
"""

import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt


ASSETS_DIR = Path(__file__).parent.parent / "assets"
TEMPLATE_PATH = ASSETS_DIR / "template.pptx"


def open_template(dest_path: str | Path | None = None):
    """Abre una copia de la plantilla. Si pasas `dest_path`, copia primero
    el archivo ahí (recomendado para no mutar la plantilla original).
    """
    if dest_path is not None:
        shutil.copy(TEMPLATE_PATH, dest_path)
        return Presentation(str(dest_path))
    return Presentation(str(TEMPLATE_PATH))


def duplicate_slide(prs, source_index: int):
    """Duplica una slide existente y la añade al final. Devuelve la nueva slide.

    source_index es base 0 (slide 2 de la plantilla → index 1).
    """
    source = prs.slides[source_index]
    blank_layout = source.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    # Copia todos los shapes del source
    # Vaciar los shapes que vienen del layout por defecto
    for shp in list(new_slide.shapes):
        shp.element.getparent().remove(shp.element)

    for shp in source.shapes:
        el = copy.deepcopy(shp.element)
        new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    return new_slide


def remove_slide(prs, index: int):
    """Elimina una slide por índice base 0."""
    slide_id_lst = prs.slides._sldIdLst
    slides = list(slide_id_lst)
    slide = slides[index]
    rId = slide.get(qn("r:id"))
    prs.part.drop_rel(rId)
    slide_id_lst.remove(slide)


def keep_only_slides(prs, indices_to_keep):
    """Deja sólo las slides cuyos índices (base 0) figuran en la lista.
    Útil cuando partes de la plantilla de 36 slides y sólo quieres 6.
    Los índices se aplican sobre el estado INICIAL.
    """
    indices_to_keep = set(indices_to_keep)
    total = len(prs.slides)
    # Borrar de atrás hacia delante para no desplazar índices
    for i in range(total - 1, -1, -1):
        if i not in indices_to_keep:
            remove_slide(prs, i)


def replace_text_in_slide(slide, replacements: dict):
    """Reemplaza texto a nivel de run para preservar formato.

    `replacements` es un dict {original: nuevo}. Cada run cuyo texto coincida
    (exacto) con una clave se sustituye. Para reemplazos parciales, el matching
    es por substring dentro del run.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def set_slide_text(slide, shape_index: int, new_text: str):
    """Sustituye el texto completo de un shape por posición, preservando el
    formato del primer run."""
    shape = slide.shapes[shape_index]
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Conservar formato del primer párrafo / run
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_run = first_para.runs[0]
        # Vaciar el resto de runs del primer párrafo
        for r in list(first_para.runs[1:]):
            r._r.getparent().remove(r._r)
        first_run.text = new_text
    else:
        first_para.text = new_text
    # Eliminar párrafos extra
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def insert_png_into_placeholder(slide, placeholder_shape, png_path: str | Path):
    """Sustituye un shape placeholder (p.ej. un rectángulo gris "INSERT PICTURE HERE")
    por una imagen PNG ocupando exactamente la misma caja.
    """
    left, top, width, height = (
        placeholder_shape.left,
        placeholder_shape.top,
        placeholder_shape.width,
        placeholder_shape.height,
    )
    # Eliminar el placeholder
    placeholder_shape.element.getparent().remove(placeholder_shape.element)
    # Insertar imagen en ese bbox
    return slide.shapes.add_picture(str(png_path), left, top, width=width, height=height)


def find_shape_by_text(slide, needle: str):
    """Devuelve el primer shape cuyo texto contiene `needle` (case-insensitive).
    Útil para localizar placeholders tipo "INSERT PICTURE HERE" o "Title goes here".
    """
    needle_lower = needle.lower()
    for shape in slide.shapes:
        if shape.has_text_frame and needle_lower in shape.text_frame.text.lower():
            return shape
    return None


def apply_font(run, *, name=None, size_pt=None, color=None, bold=None):
    """Helper para aplicar tipografía de marca a un run concreto."""
    if name is not None:
        run.font.name = name
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold
